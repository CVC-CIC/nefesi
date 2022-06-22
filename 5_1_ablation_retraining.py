"""
This file contains a toy examples to have a first contact with Nefesi, and Keras.
This file has been created with tensorflow (and tensorflow-gpu) 1.8.0, keras 2.2.0, and python 3.6 (with anaconda3 interpreter)
"""
import cv2
import torch
from torchvision import transforms
import numpy as np
import functools
from functions.network_data2 import NetworkData
from torch.utils.data import Dataset, DataLoader, random_split
BATCH_SIZE = 100
import torch.nn as nn
import matplotlib.pyplot as plt
import os
from pptx import Presentation
import pptx
from Model_generation.Unet import UNet, CAN
from datasets import FivekDataset

from torch_utils import JoinedDataLoader, load_model

def preproces_hpeu_img(imgs_hr, small_size=(300, 300)):
    imgs_lr = cv2.resize(np.array(imgs_hr), small_size)

    gray = 1 - (0.299 * (imgs_lr[:, :, 0] + 1) + 0.587 * (imgs_lr[:, :, 1] + 1) + 0.114 * (imgs_lr[:, :, 2] + 1)) / 2
    gray = [np.stack((gray,) * 3, axis=-1)]

    imgs_lr = [imgs_lr]
    imgs = imgs_lr + gray
    tnsr = [transforms.ToTensor()(img) for img in imgs]

    return tnsr


list_delete=[]
list_replace=[]


def preproces_hpeu_single_img(imgs_hr, small_size=(1000, 800)):
    imgs_lr = cv2.resize(np.array(imgs_hr), small_size)

    gray = 1 - (0.299 * (imgs_lr[:, :, 0] + 1) + 0.587 * (imgs_lr[:, :, 1] + 1) + 0.114 * (imgs_lr[:, :, 2] + 1)) / 2
    gray = [np.stack((gray,) * 3, axis=-1)]

    imgs_lr = [imgs_lr]
    imgs = imgs_lr + gray
    tnsr = [torch.unsqueeze(transforms.ToTensor()(img).float().cuda(), 0) for img in imgs]

    return tnsr


def rgetattr(obj, attr, *args):
    def _getattr(obj, attr):
        return getattr(obj, attr, *args)

    return functools.reduce(_getattr, [obj] + attr.split('.'))




def ablation_replacement_hook(idx, inputs, output):

    for i in range(len(list_delete)):

        deleted= output[:, list_delete[i], :, :]
        output[:, list_replace[i], :, :] +=deleted
        output[:, list_delete[i], :, :] = 0

    return output

def ablation_hook(idx, inputs, output):

    output[:, list_delete, :, :] = 0

    return output

def total_ablation_hook(idx, inputs, output):

    output[:, :, :, :] = 0

    return output

def main():

    image_path='/data/134-1/datasets/5K_mit_adobe/datasets'
    device=0
    file_name = '/home/guillem/Nefesi2022/nefesi/Nefesi_models/CAN32/CAN32.obj'
    Nefesimodel = NetworkData.load_from_disk(file_name, model_file=None)
    folder_dir = "/home/guillem/Nefesi2022/"
    model = CAN(n_channels=32)
    model.load_state_dict(
        torch.load(folder_dir + 'nefesi/Model_generation/Can32.pt', map_location={'cuda:1': 'cuda:0'}))

    percent=[0,10,20,30,40,50,60,70,80,90,100]
    layers = Nefesimodel.get_layer_names_to_analyze()
    prs = Presentation()
    for layer in layers:
        print(layer)
        target_layer=Nefesimodel.get_layer_by_name(layer)
        similarity=target_layer.similarity_index

        np.fill_diagonal(similarity,0)

        ablated_num=(len(similarity))
        ablated_list=np.zeros(ablated_num)
        no_ablated_list = np.zeros(ablated_num)
        for i in range(ablated_num):
            max_pos=np.unravel_index(np.argmax(similarity),similarity.shape)



            similarity_neu1=sum(similarity[max_pos[0]])
            similarity_neu2 = sum(similarity[max_pos[1]])
            if similarity_neu1 > similarity_neu2:
                ablated_neuron=max_pos[0]
                no_ablated=max_pos[1]
            else:
                ablated_neuron=max_pos[1]
                no_ablated =max_pos[0]

            ablated_list[i]=ablated_neuron
            no_ablated_list[i]=no_ablated
            similarity[ablated_neuron]=0
            similarity[:,ablated_neuron] =0

        batch_size=100
        model.eval()
        model.to(device)
        landscape_transform = transforms.Compose([
            transforms.Resize((64, 64)),
            transforms.ToTensor(),
            transforms.Normalize((0.5, 0.5, 0.5), (0.5, 0.5, 0.5))  # normalize in [-1,1]
        ])
        portrait_transform = transforms.Compose([
            transforms.Resize((64, 64)),
            transforms.ToTensor(),
            transforms.Normalize((0.5, 0.5, 0.5), (0.5, 0.5, 0.5))  # normalize in [-1,1]
        ])
        landscape_dataset = FivekDataset(image_path, expert_idx=2, transform=landscape_transform,
                                         filter_ratio="landscape")
        portrait_dataset = FivekDataset(image_path, expert_idx=2, transform=portrait_transform, filter_ratio="portrait")

        train_size = int(0.8 * len(landscape_dataset))
        test_size = len(landscape_dataset) - train_size
        train_landscape_dataset, test_landscape_dataset = random_split(landscape_dataset, [train_size, test_size])

        train_size = int(0.8 * len(portrait_dataset))
        test_size = len(portrait_dataset) - train_size
        train_portrait_dataset, test_portrait_dataset = random_split(portrait_dataset, [train_size, test_size])



        test_landscape_loader = DataLoader(landscape_dataset, batch_size=batch_size, shuffle=True, num_workers=1)
        test_portrait_loader = DataLoader(portrait_dataset, batch_size=batch_size, shuffle=True, num_workers=1)
        test_loader = JoinedDataLoader(test_landscape_loader, test_portrait_loader)

        criterion = nn.MSELoss()

        model.eval()
        global list_delete

        global list_replace

        # list_delete = list(range(0, 64))
        # model.inc.register_forward_hook(ablation_hook)
        # test_loss = []
        # for i, (im_o, im_t) in enumerate(test_loader):
        #     im_o, im_t = im_o.to(device), im_t.to(device)
        #     with torch.no_grad():
        #         output = model(im_o)
        #         test_loss.append(criterion(output, im_t).item())
        #         avg_loss = sum(test_loss) / len(test_loss)
        # print(avg_loss)
        #

        ablated_list = ablated_list.astype(int)
        no_ablated_list = no_ablated_list.astype(int)

        final_loss2 = []
        # ablated_list.sort()

        for i in range(10):
            i = i * int(ablated_num / 10)
            list_delete = ablated_list[:i]
            list_replace = no_ablated_list[:i]
            Hook = rgetattr(model,layer).register_forward_hook(ablation_replacement_hook)
            test_loss = []
            for i, (im_o, im_t) in enumerate(test_loader):
                im_o, im_t = im_o.to(device), im_t.to(device)
                with torch.no_grad():
                    output = model(im_o)
                    test_loss.append(criterion(output, im_t).item())
                    avg_loss = sum(test_loss) / len(test_loss)
            Hook.remove()
            final_loss2.append(avg_loss)


        Hook =rgetattr(model,layer).register_forward_hook(total_ablation_hook)

        test_loss = []
        for i, (im_o, im_t) in enumerate(test_loader):
            im_o, im_t = im_o.to(device), im_t.to(device)
            with torch.no_grad():
                output = model(im_o)
                test_loss.append(criterion(output, im_t).item())
                avg_loss = sum(test_loss) / len(test_loss)
        final_loss2.append(avg_loss)

        Hook.remove()




        final_loss = []


        for i in range(10):
            i = i * int(ablated_num / 10)
            list_delete = ablated_list[:i]
            Hook = rgetattr(model, layer).register_forward_hook(ablation_hook)
            test_loss = []
            for i, (im_o, im_t) in enumerate(test_loader):
                im_o, im_t = im_o.to(device), im_t.to(device)
                with torch.no_grad():
                    output = model(im_o)
                    test_loss.append(criterion(output, im_t).item())
                    avg_loss = sum(test_loss) / len(test_loss)
            Hook.remove()
            final_loss.append(avg_loss)


        Hook = rgetattr(model, layer).register_forward_hook(total_ablation_hook)

        test_loss = []
        for i, (im_o, im_t) in enumerate(test_loader):
            im_o, im_t = im_o.to(device), im_t.to(device)
            with torch.no_grad():
                output = model(im_o)
                test_loss.append(criterion(output, im_t).item())
                avg_loss = sum(test_loss) / len(test_loss)
        final_loss.append(avg_loss)
        Hook.remove()

        final_loss2=[np.abs(f-final_loss2[0]) for f in final_loss2]
        final_loss=[np.abs(f-final_loss[0]) for f in final_loss]


        # plt.subplot(2, 1, 1)
        plt.clf()
        plt.plot(percent,final_loss, label="Without reclacing weights")
        plt.plot(percent,final_loss2, 'r', label="Reclacing weights")
        plt.title('Ordered by relebance')
        plt.ylabel('Diference in loss from the OG model')
        plt.xlabel('% of neurons ablated')
        plt.legend(loc="upper left")


        # print('ordered part')
        #
        # ablated_list.sort()
        # final_loss2 = []
        #
        # for i in range(10):
        #     i = i * int(ablated_num / 10)
        #     list_delete = ablated_list[:i]
        #     list_replace = no_ablated_list[:i]
        #     Hook = getattr(model, layer).register_forward_hook(ablation_replacement_hook)
        #     test_loss = []
        #     for i, (im_o, im_t) in enumerate(test_loader):
        #         im_o, im_t = im_o.to(device), im_t.to(device)
        #         with torch.no_grad():
        #             output = model(im_o)
        #             test_loss.append(criterion(output, im_t).item())
        #             avg_loss = sum(test_loss) / len(test_loss)
        #     Hook.remove()
        #     final_loss2.append(avg_loss)
        #
        # Hook = getattr(model, layer).register_forward_hook(total_ablation_hook)
        #
        # test_loss = []
        # for i, (im_o, im_t) in enumerate(test_loader):
        #     im_o, im_t = im_o.to(device), im_t.to(device)
        #     with torch.no_grad():
        #         output = model(im_o)
        #         test_loss.append(criterion(output, im_t).item())
        #         avg_loss = sum(test_loss) / len(test_loss)
        # final_loss2.append(avg_loss)
        #
        # Hook.remove()
        #
        # final_loss = []
        #
        # for i in range(10):
        #     i = i * int(ablated_num / 10)
        #     list_delete = ablated_list[:i]
        #     Hook = getattr(model, layer).register_forward_hook(ablation_hook)
        #     test_loss = []
        #     for i, (im_o, im_t) in enumerate(test_loader):
        #         im_o, im_t = im_o.to(device), im_t.to(device)
        #         with torch.no_grad():
        #             output = model(im_o)
        #             test_loss.append(criterion(output, im_t).item())
        #             avg_loss = sum(test_loss) / len(test_loss)
        #     Hook.remove()
        #     final_loss.append(avg_loss)
        #
        # Hook = getattr(model, layer).register_forward_hook(total_ablation_hook)
        #
        # test_loss = []
        # for i, (im_o, im_t) in enumerate(test_loader):
        #     im_o, im_t = im_o.to(device), im_t.to(device)
        #     with torch.no_grad():
        #         output = model(im_o)
        #         test_loss.append(criterion(output, im_t).item())
        #         avg_loss = sum(test_loss) / len(test_loss)
        # final_loss.append(avg_loss)
        # Hook.remove()
        #
        #
        #
        #
        #
        # plt.subplot(2, 1, 2)
        #
        #
        #
        #
        # plt.plot(final_loss, label="Without reclacing weights")
        # plt.plot(final_loss2, 'r', label="Reclacing weights")
        # plt.title('Random Order')
        # plt.legend(loc="upper left")

        plt.savefig('temp_fig.jpeg')
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        slide.placeholders[0].text = layer
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        placeholder = slide.shapes
        picture = placeholder.add_picture("temp_fig.jpeg", pptx.util.Inches(0.5), pptx.util.Inches(0.5),
                                          width=pptx.util.Inches(7),
                                          height=pptx.util.Inches(7))

    prs.save("ablation study_Can32.pptx")












if __name__ == '__main__':
    main()

