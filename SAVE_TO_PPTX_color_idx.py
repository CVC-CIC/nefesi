"""
This file contains a toy examples to have a first contact with Nefesi, and Keras.
This file has been created with tensorflow (and tensorflow-gpu) 1.8.0, keras 2.2.0, and python 3.6 (with anaconda3 interpreter)
"""
import cv2
import torch
from torchvision import transforms
import numpy as np
import functools
from functions.network_data2 import NetworkData
from pptx import Presentation
import pptx
import types
import functions.GPUtil as gpu
BATCH_SIZE = 100
from HPEU.module_hpeu_gf import DeepGuidedFilterGuidedMapConvGF
from functions.image import ImageDataset
from functions.read_activations import get_activations
import interface_DeepFramework.DeepFramework as DeepF
import matplotlib.pyplot as plt
import io
from PIL import Image


def preproces_hpeu_img( imgs_hr,small_size=(64,64)):

    imgs_lr=cv2.resize(np.array(imgs_hr),small_size)

    gray = 1 - (0.299 * (imgs_lr[:, :, 0] + 1) + 0.587 * (imgs_lr[:, :, 1] + 1) + 0.114 * (imgs_lr[:, :, 2] + 1)) / 2
    gray = [np.stack((gray,) * 3, axis=-1)]

    imgs_lr = [imgs_lr]
    imgs = imgs_lr + gray
    tnsr = [transforms.ToTensor()(img) for img in imgs]


    return tnsr


def preproces_hpeu_single_img(imgs_hr, small_size=(500, 350)):
    imgs_lr = cv2.resize(np.array(imgs_hr), small_size)

    gray = 1 - (0.299 * (imgs_lr[:, :, 0] + 1) + 0.587 * (imgs_lr[:, :, 1] + 1) + 0.114 * (imgs_lr[:, :, 2] + 1)) / 2
    gray = [np.stack((gray,) * 3, axis=-1)]

    imgs_lr = [imgs_lr]
    imgs = imgs_lr + gray
    tnsr = [torch.unsqueeze(transforms.ToTensor()(img).float().cuda(), 0) for img in imgs]

    return tnsr




def main():

    ## This script creates a powerpoint with all the neuron features of your Nefesimodel sorted by its color index



    file_name= 'C:/Users/g84194584/Desktop/Code/nefesi/nefesi/hpeu_praveen/Pos_Neg.obj'

    Nefesimodel=NetworkData.load_from_disk(file_name, model_file=None)
    model_path = "C:/Users/g84194584/Desktop/Code/nefesi/HPEU/mse_net_latest.pth"
    device = torch.device('cuda:{}'.format(0))
    model = DeepGuidedFilterGuidedMapConvGF()
    model.load_state_dict(torch.load(model_path))
    model = model.lr
    model = DeepF.deep_model(model)
    Nefesimodel.model=model

    dataset=Nefesimodel.dataset

    layers=Nefesimodel.get_layer_names_to_analyze()

    prs = Presentation()



    for layer in layers:
        color_idx=[]
        nf=[]

        for n in Nefesimodel.get_layer_by_name(layer).neurons_data:
            color_idx.append(n.selectivity_idx['color'])
            nf.append(n._neuron_feature)
        sortorder=np.argsort(color_idx)[::-1]
        values=np.sort(color_idx)[::-1]
        nf_ordered=[nf[n] for n in sortorder]
        pic_size = int(nf[0].shape[0])

        num_images=12
        if len(color_idx) > 97:
            num_images = 24
        show_neurons=np.zeros((int(pic_size*len(nf_ordered)/num_images),num_images*pic_size,3))
        for i, imatge in enumerate(nf_ordered):

            res= int(pic_size*(i % (len(nf_ordered)/ num_images)))

            div = int(pic_size * (i // (len(nf_ordered) / num_images)))

            show_neurons[res:res+pic_size,div:div+pic_size]=imatge
        values[values>1]=1
        values[values <0] = 0
        plt.clf()
        plt.plot(values)
        plt.title(layer)
        plt.savefig('temp_fig.jpeg')
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        placeholder = slide.shapes
        show_neurons=show_neurons.astype(np.uint8)
        im = Image.fromarray(show_neurons)


        im=im.resize([x*10 for x in im.size], Image.NEAREST)


        im.save("temp.jpeg")

        picture = placeholder.add_picture("temp.jpeg",pptx.util.Inches(1),pptx.util.Inches(0.5),width=pptx.util.Inches(8), height=int(pptx.util.Inches(8)/num_images*(len(nf_ordered) / num_images)))
        slide.placeholders[0].text = layer
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        placeholder = slide.shapes
        picture = placeholder.add_picture("temp_fig.jpeg", pptx.util.Inches(0.5), pptx.util.Inches(0.5),
                                          width=pptx.util.Inches(7),
                                          height=pptx.util.Inches(7))

    prs.save("color_sorted_praveen.pptx")










if __name__ == '__main__':
    main()